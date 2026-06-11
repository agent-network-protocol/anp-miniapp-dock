from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from PIL import Image, ImageEnhance
try:
    from pptx.enum.dml import MSO_THEME_COLOR
except Exception:
    MSO_THEME_COLOR = None
try:
    from pptx.enum.dml import MSO_ARROWHEAD
except Exception:
    MSO_ARROWHEAD = None

ROOT = Path(__file__).resolve().parent
ASSET = ROOT / "assets"
OUT = ROOT / "anp-miniapp-dock-sharing.pptx"

W, H = 13.333333, 7.5

COLORS = {
    "bg": "071627",
    "bg2": "0B1F36",
    "panel": "102A44",
    "panel2": "123456",
    "cyan": "22D3EE",
    "teal": "2DD4BF",
    "blue": "60A5FA",
    "purple": "A78BFA",
    "amber": "F59E0B",
    "green": "34D399",
    "red": "FB7185",
    "text": "E6F1FF",
    "muted": "94A3B8",
    "line": "24445F",
    "white": "FFFFFF",
}

FONT_CN = "PingFang SC"
FONT_EN = "Aptos"
FONT_MONO = "Menlo"


def rgb(hex_str):
    hex_str = hex_str.strip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def add_bg(slide, color="bg"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(H))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[color])
    shape.line.fill.background()
    return shape


def set_transparency(fill, value):
    # python-pptx accepts percentage integer in recent versions; ignore if unavailable.
    try:
        fill.transparency = value
    except Exception:
        pass


def dim_image_path(img_path, alpha_overlay=0):
    """Return a project-local dimmed copy because python-pptx has no portable fill alpha."""
    img_path = Path(img_path)
    if not alpha_overlay:
        return img_path
    out = img_path.with_name(f"{img_path.stem}-dim{int(alpha_overlay)}{img_path.suffix}")
    if out.exists() and out.stat().st_mtime >= img_path.stat().st_mtime:
        return out
    img = Image.open(img_path).convert("RGB")
    # Keep the visual visible while making it safe behind white/cyan title text.
    brightness = max(0.18, 1.0 - float(alpha_overlay) / 100.0)
    img = ImageEnhance.Brightness(img).enhance(brightness)
    black = Image.new("RGB", img.size, (4, 14, 28))
    img = Image.blend(img, black, min(0.36, float(alpha_overlay) / 180.0))
    img.save(out)
    return out


def add_image_full(slide, img_path, alpha_overlay=0, overlay_color="bg"):
    img = dim_image_path(img_path, alpha_overlay)
    slide.shapes.add_picture(str(img), Inches(0), Inches(0), width=Inches(W), height=Inches(H))


def add_text(slide, text, x, y, w, h, size=22, color="text", bold=False, font=FONT_CN,
             align="left", valign="top", line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(valign, MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(COLORS[color])
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    if line_spacing:
        p.line_spacing = line_spacing
    return tb


def add_multiline(slide, lines, x, y, w, h, size=19, color="text", bullet=False, font=FONT_CN,
                  bold_first=False, line_spacing=1.08):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(COLORS[color])
        p.font.bold = bool(bold_first and idx == 0)
        p.line_spacing = line_spacing
        if bullet:
            p.level = 0
            p.text = f"• {line}"
    return tb


def add_title(slide, title, subtitle=None, section=None):
    if section:
        add_text(slide, section.upper(), 0.65, 0.32, 3.4, 0.25, size=10, color="cyan", bold=True, font=FONT_EN)
    add_text(slide, title, 0.62, 0.58, 10.9, 0.62, size=28, color="text", bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.66, 1.15, 11.4, 0.35, size=13, color="muted")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.66), Inches(1.52), Inches(1.25), Inches(0.035))
    line.fill.solid(); line.fill.fore_color.rgb = rgb(COLORS["cyan"]); line.line.fill.background()


def add_footer(slide, idx):
    add_text(slide, "anp-miniapp-dock · ANP DID × MiniApp MCP", 0.64, 7.14, 5.2, 0.2, size=8.5, color="muted", font=FONT_EN)
    add_text(slide, f"{idx:02d}", 12.35, 7.10, 0.45, 0.25, size=9, color="muted", font=FONT_EN, align="right")


def add_box(slide, text, x, y, w, h, fill="panel", line="line", text_color="text", size=16, bold=False,
            radius=True, font=FONT_CN, align="center", valign="middle"):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(COLORS[fill])
    if line:
        shape.line.color.rgb = rgb(COLORS[line]); shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12); tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(valign, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = rgb(COLORS[text_color])
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
    return shape


def add_chip(slide, text, x, y, w, fill="panel2", color="cyan", size=10):
    return add_box(slide, text, x, y, w, 0.32, fill=fill, line=color, text_color="text", size=size, bold=True)


def add_arrow(slide, x1, y1, x2, y2, color="cyan", width=2.0, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = rgb(COLORS[color])
    conn.line.width = Pt(width)
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    try:
        conn.line.end_arrowhead = MSO_ARROWHEAD.TRIANGLE
    except Exception:
        pass
    return conn


def add_metric(slide, label, value, x, y, w, color="cyan"):
    add_text(slide, value, x, y, w, 0.35, size=24, color=color, bold=True, font=FONT_EN, align="center")
    add_text(slide, label, x, y+0.38, w, 0.34, size=10.5, color="muted", align="center")


def add_panel(slide, x, y, w, h, title=None):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid(); panel.fill.fore_color.rgb = rgb(COLORS["bg2"]); set_transparency(panel.fill, 12)
    panel.line.color.rgb = rgb(COLORS["line"]); panel.line.width = Pt(1)
    if title:
        add_text(slide, title, x+0.22, y+0.18, w-0.44, 0.28, size=13, color="cyan", bold=True)
    return panel


def add_code_block(slide, text, x, y, w, h, size=12):
    shape = add_box(slide, "", x, y, w, h, fill="bg", line="line", size=size, align="left", valign="top", radius=True)
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = False
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.08); tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_MONO
    p.font.size = Pt(size)
    p.font.color.rgb = rgb(COLORS["text"])
    return shape


def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_image_full(slide, ASSET / "open-did-bridge.png", alpha_overlay=32)
    add_text(slide, "开源的智能体原生\n小程序 MCP 容器", 0.62, 0.78, 6.4, 1.45, size=35, color="text", bold=True)
    add_text(slide, "用 ANP DID 打通平台身份之外的开放服务入口", 0.70, 2.45, 6.6, 0.42, size=17, color="cyan", bold=True)
    add_multiline(slide, ["兼容小程序 MCP 契约", "底层身份、鉴权、网络替换为 ANP/Rust Runtime", "面向 Agent 对话中的选择、确认、支付与状态卡片"],
                  0.72, 3.05, 5.65, 1.25, size=15.5, color="text")
    add_text(slide, "15 分钟技术分享", 0.72, 6.65, 2.2, 0.25, size=10.5, color="muted")
    add_footer(slide, 1)


def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, "为什么需要新的容器？", "小程序 MCP 形态很适合 Agent，但身份入口容易被平台登录态锁住", "WHY")
    add_panel(slide, 0.62, 1.88, 4.45, 4.55, "核心矛盾")
    add_multiline(slide, [
        "小程序 MCP 已证明：原子接口 + 原子组件 + Skill 适合模型调用真实业务。",
        "但登录通常基于平台身份，例如微信登录态。",
        "商家可以支持大平台或自有账号，却很难为每个中小 Agent 单独适配。",
        "结果：Agent 想调用服务，身份和授权先成为瓶颈。",
    ], 0.85, 2.42, 3.95, 3.55, size=15.5, color="text", bullet=True)
    # identity silos diagram
    add_text(slide, "平台身份孤岛", 6.0, 1.85, 5.9, 0.32, size=16, color="muted", bold=True, align="center")
    centers = [(6.35, 3.0, "平台 A\n登录态", "blue"), (8.55, 2.65, "平台 B\n登录态", "purple"), (10.75, 3.0, "商家自有\n账号", "amber")]
    for x, y, label, col in centers:
        cyl = slide.shapes.add_shape(MSO_SHAPE.CAN, Inches(x), Inches(y), Inches(1.45), Inches(1.75))
        cyl.fill.solid(); cyl.fill.fore_color.rgb = rgb(COLORS["panel"]); cyl.line.color.rgb = rgb(COLORS[col]); cyl.line.width = Pt(2)
        cyl.text_frame.text = label
        for p in cyl.text_frame.paragraphs:
            p.font.name = FONT_CN; p.font.size = Pt(12.5); p.font.color.rgb = rgb(COLORS["text"]); p.alignment = PP_ALIGN.CENTER
    add_box(slide, "用户 Agent", 8.02, 5.75, 1.8, 0.58, fill="panel2", line="cyan", text_color="text", size=14, bold=True)
    for x, y, _, col in centers:
        add_arrow(slide, 8.92, 5.75, x+0.72, y+1.75, color="muted", width=1.2, dashed=True)
    add_box(slide, "每个平台都要单独适配\n→ 对开放 Agent 网络不友好", 6.2, 5.05, 5.8, 0.56, fill="bg2", line="red", text_color="text", size=13)
    add_footer(slide, 2)


def slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, "我们的答案：开放 DID 身份层", "不替代大厂身份，而是在其上叠加可验证、跨平台、服务方可独立校验的身份", "IDENTITY")
    # layers
    add_box(slide, "微信 / 支付宝 / 商家自有身份体系", 1.05, 5.55, 11.25, 0.58, fill="panel", line="line", size=15, bold=True)
    add_box(slide, "ANP DID 开放身份层", 1.55, 4.45, 10.25, 0.7, fill="panel2", line="cyan", size=18, bold=True)
    add_text(slide, "在平台身份之上叠加开放认证方案", 3.65, 4.06, 6.0, 0.28, size=12.5, color="muted", align="center")
    add_box(slide, "用户 / 用户 Agent\n持有 DID", 0.95, 2.35, 2.05, 0.86, fill="bg2", line="teal", size=14, bold=True)
    merchant_x = [5.05, 7.4, 9.75]
    for i, x in enumerate(merchant_x, start=1):
        add_box(slide, f"商家 Agent {i}\n校验 DID", x, 2.12 + (i%2)*0.45, 1.85, 0.84, fill="bg2", line="blue" if i != 2 else "purple", size=13, bold=True)
        add_arrow(slide, 3.0, 2.78, x, 2.55 + (i%2)*0.45, color="cyan", width=1.8)
    add_arrow(slide, 2.0, 3.22, 2.0, 4.45, color="teal", width=2.1)
    add_multiline(slide, [
        "用户拿 DID 请求服务",
        "服务方解析 DID Document + 验签",
        "授权策略由服务方独立决定",
        "商家 Agent 可服务所有开放身份 Agent",
    ], 0.96, 1.36, 11.4, 0.54, size=14.5, color="text", bullet=False)
    # Four chips
    add_chip(slide, "DID-first", 1.1, 6.42, 1.4, fill="bg2", color="cyan")
    add_chip(slide, "服务端验签", 2.7, 6.42, 1.6, fill="bg2", color="teal")
    add_chip(slide, "不破坏原平台身份", 4.55, 6.42, 2.2, fill="bg2", color="purple")
    add_chip(slide, "面向所有 Agent", 7.0, 6.42, 1.9, fill="bg2", color="amber")
    add_footer(slide, 3)


def slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, "产品定位与边界", "Agentic MiniApp Container：只做 Agent 场景里必要的小程序 MCP 子集", "SCOPE")
    add_box(slide, "Agentic MiniApp Container", 3.95, 1.75, 5.4, 0.68, fill="panel2", line="cyan", size=21, bold=True)
    add_arrow(slide, 6.65, 2.43, 6.65, 2.86, color="cyan", width=2)
    add_box(slide, "兼容 MCP 契约\n+\nANP DID 身份底座", 4.45, 2.88, 4.4, 0.9, fill="bg2", line="teal", size=16, bold=True)
    add_panel(slide, 0.82, 2.35, 3.42, 3.55, "它是什么")
    add_multiline(slide, ["面向 Agent 的 Skill 运行容器", "加载 `SKILL.md` / `mcp.json`", "执行原子接口 JS", "渲染交易型原子组件", "处理用户确认与审计"],
                  1.05, 2.92, 2.92, 2.55, size=13.5, bullet=True)
    add_panel(slide, 9.08, 2.35, 3.42, 3.55, "它不是什么")
    add_multiline(slide, ["不是完整微信小程序 Runtime", "不做复杂页面路由 / TabBar", "不复刻微信账号 / 支付 / 云开发", "不做传统 App 页面容器", "不替代 aWiki daemon / IM"],
                  9.30, 2.92, 2.92, 2.55, size=13.5, bullet=True)
    add_box(slide, "设计取舍：接口契约兼容优先，运行时能力渐进替换", 2.35, 6.25, 8.65, 0.52, fill="panel", line="amber", size=15, bold=True)
    add_footer(slide, 4)


def slide5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, "总体架构：独立 Rust Runtime", "加载、执行、渲染、授权、审计由 Runtime 编排；DID/签名/网络优先复用 ANP Rust SDK", "ARCHITECTURE")
    slide.shapes.add_picture(str(ASSET / "runtime-container.png"), Inches(7.05), Inches(1.28), width=Inches(6.0), height=Inches(3.38))
    # stack
    stack = [
        ("Rust CLI / Test Runner / Future Host", "blue"),
        ("Skill Loader + MCP Contract Validator", "cyan"),
        ("Atomic API Runtime · QuickJS-NG API VM", "teal"),
        ("MiniApp Component Runtime · Render IR", "purple"),
        ("wx Compatibility Layer · Consent · Storage", "amber"),
        ("ANP SDK Adapter · DID Auth · Signed HTTP", "green"),
        ("Demo Merchant Agent Server", "blue"),
    ]
    y = 1.72
    for text, col in stack:
        add_box(slide, text, 0.82, y, 5.7, 0.48, fill="bg2", line=col, size=12.8, bold=True)
        y += 0.6
    for i in range(len(stack)-1):
        add_arrow(slide, 3.67, 2.20 + i*0.6, 3.67, 2.31 + i*0.6, color="muted", width=1.1)
    add_panel(slide, 7.35, 4.85, 5.05, 1.55, "运行职责")
    add_multiline(slide, ["Skill 包加载与契约校验", "JS sandbox 执行原子接口", "组件运行时或 fallback 渲染", "高风险动作 consent/audit", "短期 token 与多商家会话隔离"],
                  7.56, 5.25, 4.55, 0.92, size=11.8, bullet=True)
    add_footer(slide, 5)


def slide6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, "兼容小程序 MCP 契约", "目标是让商家 Skill 尽量不重写业务代码，只替换底层身份、网络和宿主能力", "CONTRACT")
    tree = """coffee-skill/
├─ SKILL.md
├─ mcp.json
├─ index.js
├─ apis/
│  ├─ searchDrinks.js
│  ├─ confirmOrder.js
│  └─ payOrder.js
└─ components/
   ├─ drink-list/
   ├─ order-confirm/
   └─ payment-result/"""
    add_code_block(slide, tree, 0.78, 1.88, 4.0, 3.65, size=11.5)
    add_text(slide, "Skill 包结构保持原样", 1.0, 5.68, 3.6, 0.28, size=13, color="muted", align="center")
    chips = [
        ("`mcp.json.apis[]`", 5.35, 1.9, 2.05, "cyan"),
        ("`inputSchema`", 7.58, 1.9, 1.62, "teal"),
        ("`outputSchema`", 9.42, 1.9, 1.78, "purple"),
        ("`structuredContent`", 5.35, 2.38, 2.35, "blue"),
        ("`_meta.ui.componentPath`", 7.93, 2.38, 2.8, "amber"),
        ("`sendFollowUpMessage`", 5.35, 2.86, 2.65, "green"),
        ("`api/call`", 8.25, 2.86, 1.45, "cyan"),
        ("middleware", 9.95, 2.86, 1.45, "purple"),
    ]
    for txt, x, y, w, col in chips:
        add_chip(slide, txt, x, y, w, fill="bg2", color=col, size=9.2)
    add_panel(slide, 5.25, 3.55, 6.95, 1.75, "返回结构的语义")
    add_multiline(slide, [
        "`content`：给模型/Agent 的文本引导",
        "`structuredContent`：模型可见，也作为组件渲染数据",
        "`_meta`：模型不可见，传递组件私有数据",
        "组件事件必须回到统一 Orchestrator",
    ], 5.50, 4.0, 6.35, 0.95, size=12.8, bullet=True)
    add_box(slide, "兼容 ≠ 完整复刻；契约保持，底层替换", 5.45, 5.78, 6.45, 0.55, fill="panel", line="cyan", size=15, bold=True)
    add_footer(slide, 6)


def slide7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, "组件运行时：从 WXML/WXSS 到 Render IR", "组件主线不是纯 CardSpec，而是 MiniApp MCP 原子组件运行时子集", "COMPONENTS")
    boxes = [
        ("componentPath\n组件目录", 0.85, "cyan"),
        ("Component JS\nWXML/WXSS 子集", 3.05, "teal"),
        ("QuickJS-NG\nComponent VM", 5.45, "purple"),
        ("Render IR\n平台无关树", 7.82, "blue"),
        ("Conversation\nCard", 10.15, "amber"),
    ]
    for label, x, col in boxes:
        add_box(slide, label, x, 2.15, 1.85, 0.9, fill="bg2", line=col, size=13.2, bold=True)
    for i in range(len(boxes)-1):
        add_arrow(slide, boxes[i][1]+1.85, 2.60, boxes[i+1][1], 2.60, color="cyan", width=2)
    add_panel(slide, 0.9, 3.75, 5.6, 1.45, "P0 支持子集")
    add_multiline(slide, ["`Component({})`、data/properties/methods、生命周期", "`setData`、`bindtap`、Input/Result/Expire", "view/text/image/button/scroll-view、wx:if/wx:for、简单绑定"],
                  1.13, 4.15, 5.15, 0.7, size=11.4, bullet=True)
    add_panel(slide, 7.0, 3.75, 5.15, 1.45, "Fallback 梯子")
    ladder = [("1 Component Runtime", "cyan"), ("2 native adapter", "teal"), ("3 CardSpec", "amber"), ("4 content text", "muted")]
    x0 = 7.28
    for i, (label, col) in enumerate(ladder):
        add_box(slide, label, x0 + i*1.15, 4.25 + i*0.18, 1.1, 0.36, fill="bg2", line=col if col != "muted" else "line", size=8.7, bold=True)
    add_box(slide, "关键原则：WXML AST 不直接绑定宿主 UI，Render IR 隔离模板语义和渲染后端", 1.15, 6.0, 10.95, 0.52, fill="panel", line="cyan", size=14.2, bold=True)
    add_footer(slide, 7)


def slide8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_image_full(slide, ASSET / "coffee-did-flow.png", alpha_overlay=27)
    add_title(slide, "登录流程：咖啡点单中的 DID 认证", "challenge → DID 签名 → 服务端验签 → capability token", "AUTH FLOW")
    add_panel(slide, 0.8, 1.85, 4.25, 4.65, "模拟流程")
    steps = [
        "1 发现咖啡商家 Agent",
        "2 读取 `SKILL.md` 与 `mcp.json`",
        "3 请求 challenge / 登录请求",
        "4 用户 DID 使用 ANP Rust SDK 签名",
        "5 服务端解析 DID Document 并验签",
        "6 创建/查找 DID 绑定账户",
        "7 返回短期 capability token",
    ]
    y = 2.32
    for s in steps:
        add_box(slide, s, 1.05, y, 3.75, 0.34, fill="bg2", line="line", size=10.7, align="left")
        y += 0.52
    add_panel(slide, 8.35, 4.72, 3.9, 1.23, "与 DID-WBA 一致")
    add_multiline(slide, ["`Signature-Input` / `Signature`", "body 绑定 `Content-Digest`", "检查 `authentication` 授权、公钥、时间窗口、nonce"],
                  8.55, 5.08, 3.45, 0.55, size=10.6, bullet=True)
    add_footer(slide, 8)


def slide9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, "咖啡点单闭环", "用交易 demo 验证加载 Skill、接口调用、组件渲染、用户确认、支付和卡片过期", "DEMO")
    flow = [
        ("searchDrinks", "API\n搜索饮品", "cyan"),
        ("drink-list", "组件\n选择饮品", "teal"),
        ("confirmOrder", "API\n确认订单", "purple"),
        ("order-confirm", "组件\n支付确认", "amber"),
        ("payOrder", "API\n模拟支付", "green"),
        ("payment-result", "组件\n状态展示", "blue"),
    ]
    xs = [0.75, 2.75, 4.75, 6.75, 8.75, 10.75]
    for (name, desc, col), x in zip(flow, xs):
        add_box(slide, name, x, 2.25, 1.55, 0.45, fill="bg2", line=col, size=10.5, bold=True)
        add_box(slide, desc, x, 2.82, 1.55, 0.82, fill="panel", line="line", size=10.2)
    for i in range(len(xs)-1):
        add_arrow(slide, xs[i]+1.55, 3.05, xs[i+1], 3.05, color="cyan", width=1.7)
    # consent gate
    add_box(slide, "Consent / Audit\n高风险动作边界", 6.55, 4.25, 2.0, 0.75, fill="bg2", line="red", size=12.5, bold=True)
    add_arrow(slide, 7.75, 3.65, 7.55, 4.25, color="red", width=1.7)
    add_panel(slide, 1.15, 5.32, 10.9, 0.85, "验证能力")
    add_multiline(slide, ["`api/call` 回到 Orchestrator", "组件事件触发 methods", "human authorization", "expirePreviousCards", "CardSpec/content fallback"],
                  1.37, 5.72, 10.5, 0.28, size=12.4, color="text")
    add_footer(slide, 9)


def slide10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, "安全边界与风控", "DID 证明“谁在请求”；高风险动作仍必须有用户确认、策略和审计", "SECURITY")
    cards = [
        ("JS Sandbox", ["禁止宿主全局对象", "禁止任意 FS/网络", "禁止远程代码加载", "受限 CommonJS"], "cyan"),
        ("Host Capability", ["request allowlist", "网络/存储/支付边界", "设备能力按需开放", "错误与日志捕获"], "teal"),
        ("Identity & Token", ["DID Document + 验签", "短期 capability token", "按用户/商家/Skill 隔离", "过期与重放防护"], "purple"),
        ("Consent & Audit", ["下单/支付需确认", "地址/手机号需授权", "审计高风险动作", "不让 Skill 直通执行"], "amber"),
    ]
    x = 0.7
    for title, bullets, col in cards:
        add_panel(slide, x, 2.0, 2.85, 3.6, title)
        add_box(slide, "", x+0.22, 2.42, 0.32, 0.32, fill=col if col in COLORS else "cyan", line=col if col in COLORS else "cyan")
        add_multiline(slide, bullets, x+0.30, 2.95, 2.35, 1.95, size=11.2, bullet=True)
        x += 3.05
    add_box(slide, "Consent First：下单、支付、地址、手机号、身份绑定等高风险动作必须用户确认和审计", 1.2, 6.18, 10.8, 0.55, fill="panel", line="red", size=14.5, bold=True)
    add_footer(slide, 10)


def slide11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); add_title(slide, "当前原型与后续路线", "保持 Agent-native 边界：扩大兼容面，但不走向完整微信小程序 Runtime", "ROADMAP")
    phases = [
        ("P0 · 当前 MVP", ["Skill Loader + 原子接口主链路", "QuickJS API VM / Component VM", "Render IR JSON + CardSpec fallback", "DID 登录 + token + coffee demo"], "cyan"),
        ("P1 · 能力补全", ["图片/文件、扫码、电话", "真实地址/手机号/Payment Intent", "WebSocket 子集", "openDetailPage fallback、动态组件"], "amber"),
        ("P2 · 兼容与性能", ["更多 WXML/WXSS 表达式", "交易型 Skill 测试集", "Flutter/Web Renderer Adapter", "性能、稳定性、安全审计"], "purple"),
    ]
    x = 0.95
    for title, bullets, col in phases:
        add_box(slide, title, x, 2.05, 3.45, 0.55, fill="panel2", line=col, size=15.2, bold=True)
        add_panel(slide, x, 2.78, 3.45, 2.55)
        add_multiline(slide, bullets, x+0.18, 3.1, 3.05, 1.55, size=11.8, bullet=True)
        x += 4.0
    add_arrow(slide, 4.4, 2.33, 4.95, 2.33, color="cyan", width=2.2)
    add_arrow(slide, 8.4, 2.33, 8.95, 2.33, color="cyan", width=2.2)
    add_box(slide, "坚持：MCP Interface Compatibility + ANP SDK First + Sandbox First + Consent First", 1.15, 6.15, 11.0, 0.55, fill="panel", line="cyan", size=14, bold=True)
    add_footer(slide, 11)


def slide12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_image_full(slide, ASSET / "open-did-bridge.png", alpha_overlay=42)
    add_text(slide, "一起把商家 Skill\n带到开放 Agent 网络", 0.68, 0.78, 6.1, 1.25, size=31, color="text", bold=True)
    add_text(slide, "让每个商家的智能体，都能服务所有开放身份的智能体", 0.75, 2.24, 7.0, 0.35, size=15.5, color="cyan", bold=True)
    add_panel(slide, 0.72, 3.0, 5.7, 2.65, "欢迎贡献")
    add_multiline(slide, ["更多 MiniApp MCP 兼容测试", "更多交易型组件与 Skill 示例", "host capability adapter", "Flutter / Web Renderer", "安全审计、性能优化、真实商家集成"],
                  0.98, 3.46, 5.2, 1.55, size=13.7, bullet=True)
    add_box(slide, "Open identity · Open Skill runtime · Open Agent services", 0.82, 6.24, 5.9, 0.42, fill="bg2", line="cyan", text_color="text", size=12.5, bold=True, font=FONT_EN)
    add_footer(slide, 12)


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    for fn in [slide1, slide2, slide3, slide4, slide5, slide6, slide7, slide8, slide9, slide10, slide11, slide12]:
        fn(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
